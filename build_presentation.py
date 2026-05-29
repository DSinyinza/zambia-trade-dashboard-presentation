"""
build_presentation.py
=====================
Generates a professional 8-slide PowerPoint presentation:
  "Zambia's Trade Story: K33 Billion Warning"
  Storytelling with Data principles applied.

Requirements:
    pip install python-pptx

Run:
    python build_presentation.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ─────────────────────────────────────────────
# COLOR PALETTE
# ─────────────────────────────────────────────
BG_DARK        = RGBColor(0x0d, 0x1b, 0x2a)   # Deep navy background
BG_PANEL       = RGBColor(0x16, 0x21, 0x3e)   # Slightly lighter panel
ACCENT_IMPORT  = RGBColor(0xe8, 0x5d, 0x04)   # Orange — imports / warning
ACCENT_EXPORT  = RGBColor(0x1d, 0x6f, 0xa4)   # Teal-blue — exports / positive
COLOR_DEFICIT  = RGBColor(0xc1, 0x12, 0x1f)   # Red — deficit / danger
COLOR_WHITE    = RGBColor(0xff, 0xff, 0xff)
COLOR_LGREY    = RGBColor(0xe0, 0xe0, 0xe0)
COLOR_MGREY    = RGBColor(0xa0, 0xa8, 0xb8)
COLOR_GOLD     = RGBColor(0xf4, 0xa2, 0x61)   # Warm gold for callouts
COLOR_GREEN    = RGBColor(0x2d, 0xc6, 0x53)   # Positive surplus

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ─────────────────────────────────────────────
# HELPER UTILITIES
# ─────────────────────────────────────────────

def add_rect(slide, left, top, width, height, fill_color, line_color=None, line_width=None):
    """Add a filled rectangle shape."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, text, left, top, width, height,
                 font_name="Calibri", font_size=18, bold=False, italic=False,
                 color=COLOR_WHITE, align=PP_ALIGN.LEFT, word_wrap=True,
                 v_anchor=None):
    """Add a text box with given formatting."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    if v_anchor:
        tf.vertical_anchor = v_anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_multiline_text(slide, lines, left, top, width, height,
                       font_name="Calibri", font_size=16, bold=False,
                       color=COLOR_LGREY, align=PP_ALIGN.LEFT, spacing=1.15):
    """Add a text box with multiple lines (list of strings)."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    from pptx.util import Pt as _Pt
    from pptx.oxml.ns import qn
    import lxml.etree as etree

    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font_name
        run.font.size = _Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
        # line spacing
        from pptx.oxml.ns import qn as _qn
        pPr = p._pPr
        if pPr is None:
            pPr = p._p.get_or_add_pPr()
        lnSpc = etree.SubElement(pPr, _qn('a:lnSpc'))
        spcPct = etree.SubElement(lnSpc, _qn('a:spcPct'))
        spcPct.set('val', str(int(spacing * 100000)))
    return txBox


def set_slide_background(slide, color):
    """Fill slide background with a solid color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_action_headline(slide, text, top_offset=Inches(0.25)):
    """Add the prominent action headline at the top of a content slide."""
    add_rect(slide,
             Inches(0), top_offset,
             Inches(13.33), Inches(1.0),
             BG_PANEL)
    add_text_box(slide, text,
                 Inches(0.4), top_offset + Inches(0.08),
                 Inches(12.5), Inches(0.85),
                 font_size=26, bold=True,
                 color=COLOR_WHITE,
                 align=PP_ALIGN.LEFT)


def add_slide_number(slide, number, total=8):
    add_text_box(slide, f"{number} / {total}",
                 Inches(12.0), Inches(7.1),
                 Inches(1.2), Inches(0.3),
                 font_size=9, color=COLOR_MGREY,
                 align=PP_ALIGN.RIGHT)


def add_footer(slide, text="Data: Zambia Agriculture Import-Export Dashboard 2013–2022  |  FSIO Cohort Analysis"):
    add_text_box(slide, text,
                 Inches(0.3), Inches(7.15),
                 Inches(11.5), Inches(0.28),
                 font_size=8, color=COLOR_MGREY,
                 align=PP_ALIGN.LEFT)


# ─────────────────────────────────────────────
# DATA
# ─────────────────────────────────────────────
YEARS = [2013, 2014, 2015, 2016, 2017, 2018, 2019, 2020, 2021, 2022]
EXPORTS = [5.44, 4.73, 5.91, 6.23, 5.33, 5.59, 6.59, 9.70, 14.17, 12.95]
IMPORTS = [4.21, 4.62, 7.03, 7.29, 9.43, 9.24, 10.40, 15.46, 20.13, 22.22]
BALANCE = [e - i for e, i in zip(EXPORTS, IMPORTS)]


# ─────────────────────────────────────────────
# SLIDE 1: TITLE SLIDE
# ─────────────────────────────────────────────
def build_slide_01(prs):
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, BG_DARK)

    # Left accent bar
    add_rect(slide, Inches(0), Inches(0), Inches(0.12), Inches(7.5), ACCENT_IMPORT)

    # Top decorative band
    add_rect(slide, Inches(0.12), Inches(0), Inches(13.21), Inches(0.06), ACCENT_EXPORT)

    # Subtitle label pill
    add_rect(slide, Inches(0.5), Inches(1.1), Inches(3.4), Inches(0.38), ACCENT_EXPORT)
    add_text_box(slide, "AGRICULTURE TRADE ANALYSIS  2013 – 2022",
                 Inches(0.55), Inches(1.12), Inches(3.3), Inches(0.34),
                 font_size=8.5, bold=True, color=COLOR_WHITE, align=PP_ALIGN.LEFT)

    # Main title
    add_text_box(slide, "Zambia's Trade Story:",
                 Inches(0.5), Inches(1.65), Inches(11.0), Inches(1.0),
                 font_size=52, bold=True, color=COLOR_WHITE, align=PP_ALIGN.LEFT)
    add_text_box(slide, "K33 Billion Warning",
                 Inches(0.5), Inches(2.55), Inches(11.0), Inches(1.1),
                 font_size=60, bold=True, color=ACCENT_IMPORT, align=PP_ALIGN.LEFT)

    # Subtitle
    add_text_box(slide,
                 "How Currency Volatility is Amplifying the Agriculture\nImport-Export Deficit (2013–2022)",
                 Inches(0.5), Inches(3.75), Inches(10.0), Inches(1.0),
                 font_size=19, bold=False, color=COLOR_LGREY, align=PP_ALIGN.LEFT)

    # Divider line
    add_rect(slide, Inches(0.5), Inches(4.85), Inches(5.0), Inches(0.04), COLOR_MGREY)

    # Presenter & date
    add_text_box(slide, "FSIO Cohort  |  May 2026",
                 Inches(0.5), Inches(4.97), Inches(8.0), Inches(0.4),
                 font_size=13, color=COLOR_MGREY, align=PP_ALIGN.LEFT)

    # Right side — big stat block
    add_rect(slide, Inches(9.0), Inches(1.8), Inches(3.8), Inches(4.5), BG_PANEL)
    add_text_box(slide, "CUMULATIVE DEFICIT",
                 Inches(9.15), Inches(1.98), Inches(3.5), Inches(0.35),
                 font_size=9, bold=True, color=COLOR_MGREY, align=PP_ALIGN.CENTER)
    add_text_box(slide, "-K33.4bn",
                 Inches(9.0), Inches(2.3), Inches(3.8), Inches(1.1),
                 font_size=54, bold=True, color=COLOR_DEFICIT, align=PP_ALIGN.CENTER)

    add_rect(slide, Inches(9.3), Inches(3.55), Inches(3.2), Inches(0.03), ACCENT_EXPORT)

    add_text_box(slide, "TOTAL IMPORTS",
                 Inches(9.15), Inches(3.7), Inches(3.5), Inches(0.3),
                 font_size=8.5, bold=True, color=COLOR_MGREY, align=PP_ALIGN.CENTER)
    add_text_box(slide, "K110.04bn",
                 Inches(9.0), Inches(3.95), Inches(3.8), Inches(0.65),
                 font_size=30, bold=True, color=ACCENT_IMPORT, align=PP_ALIGN.CENTER)

    add_text_box(slide, "TOTAL EXPORTS",
                 Inches(9.15), Inches(4.7), Inches(3.5), Inches(0.3),
                 font_size=8.5, bold=True, color=COLOR_MGREY, align=PP_ALIGN.CENTER)
    add_text_box(slide, "K76.65bn",
                 Inches(9.0), Inches(4.95), Inches(3.8), Inches(0.65),
                 font_size=30, bold=True, color=ACCENT_EXPORT, align=PP_ALIGN.CENTER)

    # Bottom note
    add_text_box(slide,
                 "A Storytelling with Data Analysis",
                 Inches(0.5), Inches(6.85), Inches(8.0), Inches(0.4),
                 font_size=9, italic=True, color=COLOR_MGREY)

    return slide


# ─────────────────────────────────────────────
# SLIDE 2: THE SITUATION
# ─────────────────────────────────────────────
def build_slide_02(prs):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, BG_DARK)
    add_rect(slide, Inches(0), Inches(0), Inches(0.08), Inches(7.5), ACCENT_EXPORT)

    add_action_headline(slide,
        "Zambia Flipped from Surplus to Structural Deficit in Under a Decade")

    # Left narrative
    add_text_box(slide, "THE SITUATION",
                 Inches(0.3), Inches(1.45), Inches(3.0), Inches(0.3),
                 font_size=8.5, bold=True, color=ACCENT_EXPORT)

    body_lines = [
        "In 2013, Zambia exported MORE than it imported.",
        "",
        "By 2022, the annual trade deficit reached K9.3 billion —",
        "part of a cumulative K33.4 billion shortfall over the decade.",
        "",
        "This is not a temporary shock. It is a structural shift",
        "driven by currency depreciation, import dependency,",
        "and low-value commodity exports.",
    ]
    add_multiline_text(slide, body_lines,
                       Inches(0.3), Inches(1.8), Inches(5.8), Inches(3.5),
                       font_size=15.5, color=COLOR_LGREY)

    # Callout box: 2013 surplus
    add_rect(slide, Inches(0.3), Inches(5.35), Inches(2.65), Inches(1.5), BG_PANEL)
    add_text_box(slide, "2013",
                 Inches(0.35), Inches(5.42), Inches(2.5), Inches(0.35),
                 font_size=10, bold=True, color=COLOR_MGREY)
    add_text_box(slide, "+K1.23bn",
                 Inches(0.35), Inches(5.72), Inches(2.5), Inches(0.6),
                 font_size=28, bold=True, color=COLOR_GREEN)
    add_text_box(slide, "TRADE SURPLUS",
                 Inches(0.35), Inches(6.28), Inches(2.5), Inches(0.28),
                 font_size=8.5, bold=True, color=COLOR_MGREY)

    # Callout box: 2022 deficit
    add_rect(slide, Inches(3.1), Inches(5.35), Inches(2.65), Inches(1.5), BG_PANEL)
    add_text_box(slide, "2022",
                 Inches(3.15), Inches(5.42), Inches(2.5), Inches(0.35),
                 font_size=10, bold=True, color=COLOR_MGREY)
    add_text_box(slide, "-K9.27bn",
                 Inches(3.15), Inches(5.72), Inches(2.5), Inches(0.6),
                 font_size=28, bold=True, color=COLOR_DEFICIT)
    add_text_box(slide, "TRADE DEFICIT",
                 Inches(3.15), Inches(6.28), Inches(2.5), Inches(0.28),
                 font_size=8.5, bold=True, color=COLOR_MGREY)

    # Arrow between boxes
    add_text_box(slide, "→",
                 Inches(2.78), Inches(5.85), Inches(0.4), Inches(0.5),
                 font_size=26, bold=True, color=COLOR_DEFICIT)

    # Right panel — big numbers
    add_rect(slide, Inches(6.5), Inches(1.4), Inches(6.5), Inches(5.7), BG_PANEL)
    add_text_box(slide, "10-YEAR SCORECARD",
                 Inches(6.7), Inches(1.55), Inches(6.0), Inches(0.32),
                 font_size=9.5, bold=True, color=COLOR_MGREY)

    # Imports block
    add_rect(slide, Inches(6.7), Inches(1.95), Inches(5.9), Inches(1.4), ACCENT_IMPORT)
    add_text_box(slide, "TOTAL IMPORTS",
                 Inches(6.85), Inches(2.0), Inches(5.5), Inches(0.32),
                 font_size=9, bold=True, color=COLOR_WHITE)
    add_text_box(slide, "K110.04 billion",
                 Inches(6.85), Inches(2.28), Inches(5.5), Inches(0.85),
                 font_size=38, bold=True, color=COLOR_WHITE)

    # Exports block
    add_rect(slide, Inches(6.7), Inches(3.45), Inches(5.9), Inches(1.4), ACCENT_EXPORT)
    add_text_box(slide, "TOTAL EXPORTS",
                 Inches(6.85), Inches(3.5), Inches(5.5), Inches(0.32),
                 font_size=9, bold=True, color=COLOR_WHITE)
    add_text_box(slide, "K76.65 billion",
                 Inches(6.85), Inches(3.78), Inches(5.5), Inches(0.85),
                 font_size=38, bold=True, color=COLOR_WHITE)

    # Deficit
    add_rect(slide, Inches(6.7), Inches(4.95), Inches(5.9), Inches(1.0), COLOR_DEFICIT)
    add_text_box(slide, "CUMULATIVE DEFICIT",
                 Inches(6.85), Inches(4.98), Inches(5.5), Inches(0.28),
                 font_size=9, bold=True, color=COLOR_WHITE)
    add_text_box(slide, "-K33.39 billion",
                 Inches(6.85), Inches(5.22), Inches(5.5), Inches(0.65),
                 font_size=34, bold=True, color=COLOR_WHITE)

    add_footer(slide)
    add_slide_number(slide, 2)
    return slide


# ─────────────────────────────────────────────
# SLIDE 3: DIVERGING TRENDS
# ─────────────────────────────────────────────
def build_slide_03(prs):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, BG_DARK)
    add_rect(slide, Inches(0), Inches(0), Inches(0.08), Inches(7.5), ACCENT_IMPORT)

    add_action_headline(slide,
        "Imports Grew 428%. Exports Grew 138%. The Gap Tells the Story.")

    # Chart area background
    chart_left = Inches(0.4)
    chart_top  = Inches(1.42)
    chart_w    = Inches(12.5)
    chart_h    = Inches(5.3)
    add_rect(slide, chart_left, chart_top, chart_w, chart_h, BG_PANEL)

    # Y-axis label
    add_text_box(slide, "ZMW (billions)",
                 Inches(0.42), Inches(1.5), Inches(1.2), Inches(0.28),
                 font_size=8, color=COLOR_MGREY)

    # Horizontal gridlines
    max_val = 24.0
    for grid_val in [5, 10, 15, 20]:
        y_frac = 1.0 - grid_val / max_val
        gy = chart_top + Inches(0.45) + y_frac * Inches(4.55)
        add_rect(slide, chart_left + Inches(0.6), gy,
                 chart_w - Inches(0.75), Inches(0.01),
                 RGBColor(0x2a, 0x3a, 0x52))
        add_text_box(slide, f"{grid_val}",
                     chart_left + Inches(0.05), gy - Inches(0.12),
                     Inches(0.5), Inches(0.28),
                     font_size=7.5, color=COLOR_MGREY, align=PP_ALIGN.RIGHT)

    # Bar parameters
    bar_area_left  = chart_left + Inches(0.65)
    bar_area_w     = chart_w - Inches(0.8)
    bar_area_top   = chart_top + Inches(0.45)
    bar_area_h     = Inches(4.55)
    n = len(YEARS)
    slot_w = bar_area_w / n
    bar_w  = slot_w * 0.38
    gap    = slot_w * 0.04

    for i, year in enumerate(YEARS):
        imp = IMPORTS[i]
        exp = EXPORTS[i]
        x_center = bar_area_left + (i + 0.5) * slot_w

        # Import bar
        ih = (imp / max_val) * bar_area_h
        iy = bar_area_top + bar_area_h - ih
        add_rect(slide, x_center - bar_w - gap, iy, bar_w, ih, ACCENT_IMPORT)

        # Export bar
        eh = (exp / max_val) * bar_area_h
        ey = bar_area_top + bar_area_h - eh
        add_rect(slide, x_center + gap, ey, bar_w, eh, ACCENT_EXPORT)

        # Year label
        add_text_box(slide, str(year),
                     x_center - Inches(0.32), bar_area_top + bar_area_h + Inches(0.04),
                     Inches(0.64), Inches(0.24),
                     font_size=7.5, color=COLOR_MGREY, align=PP_ALIGN.CENTER)

    # Annotations
    # 2013 surplus arrow
    add_text_box(slide, "2013: Surplus\n(Exports > Imports)",
                 Inches(0.6), Inches(2.1), Inches(1.7), Inches(0.6),
                 font_size=8, color=COLOR_GREEN)

    # 2017 deficit starts
    add_text_box(slide, "2015–17:\nDeficit\nbegan",
                 Inches(3.55), Inches(1.75), Inches(1.3), Inches(0.7),
                 font_size=8, color=ACCENT_IMPORT)

    # 2020 acceleration
    add_text_box(slide, "2020–21:\nCOVID-era\nacceleration",
                 Inches(8.8), Inches(1.6), Inches(1.7), Inches(0.68),
                 font_size=8, color=ACCENT_IMPORT)

    # Legend
    add_rect(slide, Inches(9.8), Inches(6.88), Inches(0.22), Inches(0.22), ACCENT_IMPORT)
    add_text_box(slide, "Imports",
                 Inches(10.06), Inches(6.87), Inches(1.0), Inches(0.25),
                 font_size=9, color=COLOR_LGREY)
    add_rect(slide, Inches(11.2), Inches(6.88), Inches(0.22), Inches(0.22), ACCENT_EXPORT)
    add_text_box(slide, "Exports",
                 Inches(11.46), Inches(6.87), Inches(1.0), Inches(0.25),
                 font_size=9, color=COLOR_LGREY)

    # CAGR callout
    add_text_box(slide, "Import CAGR ~18%  vs  Export CAGR ~10%",
                 Inches(0.45), Inches(6.9), Inches(7.5), Inches(0.28),
                 font_size=9.5, bold=True, color=COLOR_GOLD)

    add_footer(slide)
    add_slide_number(slide, 3)
    return slide


# ─────────────────────────────────────────────
# SLIDE 4: CURRENCY TWIST
# ─────────────────────────────────────────────
def build_slide_04(prs):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, BG_DARK)
    add_rect(slide, Inches(0), Inches(0), Inches(0.08), Inches(7.5), COLOR_GOLD)

    add_action_headline(slide,
        "The Numbers Lie: Currency Depreciation is Inflating Both Sides")

    # Label
    add_text_box(slide, "THE TWIST — CURRENCY LENS",
                 Inches(0.3), Inches(1.45), Inches(6.0), Inches(0.28),
                 font_size=8.5, bold=True, color=COLOR_GOLD)

    # ZMW/USD timeline
    add_rect(slide, Inches(0.3), Inches(1.82), Inches(5.9), Inches(3.5), BG_PANEL)
    add_text_box(slide, "ZMW / USD EXCHANGE RATE TRAJECTORY",
                 Inches(0.45), Inches(1.92), Inches(5.6), Inches(0.3),
                 font_size=8.5, bold=True, color=COLOR_MGREY)

    # Simple timeline bar showing depreciation
    timeline_left = Inches(0.5)
    timeline_top  = Inches(2.35)
    timeline_w    = Inches(5.5)
    timeline_h    = Inches(0.55)

    # Gradient simulation: narrow green → wide red
    segments = [
        (0.0,  0.08, COLOR_GREEN),
        (0.08, 0.22, RGBColor(0x8b, 0xc3, 0x4a)),
        (0.22, 0.42, COLOR_GOLD),
        (0.42, 0.65, ACCENT_IMPORT),
        (0.65, 1.0,  COLOR_DEFICIT),
    ]
    for s_start, s_end, col in segments:
        add_rect(slide,
                 timeline_left + s_start * timeline_w,
                 timeline_top,
                 (s_end - s_start) * timeline_w,
                 timeline_h,
                 col)

    # Year markers
    rate_data = [(2013, "~5"), (2015, "~8"), (2017, "~9.5"),
                 (2019, "~12"), (2020, "~21"), (2022, "~17-21")]
    for year, rate in rate_data:
        frac = (year - 2013) / 9.0
        xp = timeline_left + frac * timeline_w
        add_rect(slide, xp, timeline_top + timeline_h,
                 Inches(0.02), Inches(0.15), COLOR_MGREY)
        add_text_box(slide, str(year),
                     xp - Inches(0.2), timeline_top + timeline_h + Inches(0.16),
                     Inches(0.42), Inches(0.22),
                     font_size=7.5, color=COLOR_MGREY, align=PP_ALIGN.CENTER)
        add_text_box(slide, rate,
                     xp - Inches(0.2), timeline_top - Inches(0.28),
                     Inches(0.42), Inches(0.22),
                     font_size=8, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)

    add_text_box(slide, "ZMW per USD",
                 Inches(0.5), timeline_top - Inches(0.45), Inches(1.2), Inches(0.22),
                 font_size=7.5, bold=True, color=COLOR_MGREY)

    # Key insight boxes
    insights = [
        ("IMPORT INFLATION", ACCENT_IMPORT,
         "K4.2bn (2013) → K22.2bn (2022)\nSame USD volume, 4× the ZMW cost\ndue to ~75% Kwacha depreciation"),
        ("EXPORT ILLUSION", ACCENT_EXPORT,
         "Tobacco/maize volumes barely grew.\nZMW gains are translation effects,\nnot real volume expansion"),
        ("THE REALITY", COLOR_DEFICIT,
         "\"Nominal growth is NOT real growth\nwhen your currency has lost\n75% of its value\""),
    ]
    for idx, (title, color, body) in enumerate(insights):
        bx = Inches(0.3) + idx * Inches(1.96)
        by = Inches(3.48)
        add_rect(slide, bx, by, Inches(1.86), Inches(1.72), BG_PANEL)
        add_rect(slide, bx, by, Inches(1.86), Inches(0.28), color)
        add_text_box(slide, title,
                     bx + Inches(0.05), by + Inches(0.03),
                     Inches(1.76), Inches(0.24),
                     font_size=7.5, bold=True, color=COLOR_WHITE)
        add_text_box(slide, body,
                     bx + Inches(0.07), by + Inches(0.35),
                     Inches(1.72), Inches(1.3),
                     font_size=9, color=COLOR_LGREY)

    # Right big callout
    add_rect(slide, Inches(6.5), Inches(1.42), Inches(6.5), Inches(5.7), BG_PANEL)
    add_text_box(slide, "IMPORT COST REALITY CHECK",
                 Inches(6.7), Inches(1.55), Inches(6.0), Inches(0.3),
                 font_size=9, bold=True, color=COLOR_MGREY)

    add_text_box(slide, "K4.2bn",
                 Inches(6.7), Inches(1.95), Inches(3.5), Inches(0.8),
                 font_size=42, bold=True, color=COLOR_GREEN)
    add_text_box(slide, "2013 imports at ~5 ZMW/USD",
                 Inches(6.7), Inches(2.72), Inches(5.5), Inches(0.28),
                 font_size=10, color=COLOR_MGREY)

    add_text_box(slide, "≈ same USD volume",
                 Inches(6.7), Inches(3.1), Inches(5.5), Inches(0.35),
                 font_size=13, italic=True, color=COLOR_GOLD)

    add_text_box(slide, "K22.2bn",
                 Inches(6.7), Inches(3.52), Inches(3.5), Inches(0.8),
                 font_size=42, bold=True, color=COLOR_DEFICIT)
    add_text_box(slide, "2022 imports at ~17–21 ZMW/USD",
                 Inches(6.7), Inches(4.28), Inches(5.5), Inches(0.28),
                 font_size=10, color=COLOR_MGREY)

    add_rect(slide, Inches(6.7), Inches(4.72), Inches(5.9), Inches(0.03), COLOR_MGREY)

    add_text_box(slide, "428% more Kwacha\nfor the same goods",
                 Inches(6.7), Inches(4.85), Inches(5.9), Inches(0.8),
                 font_size=21, bold=True, color=ACCENT_IMPORT)

    add_text_box(slide,
                 "ZMW weakened from ~5/USD (2013)\nto ~17–21/USD (2020–22) — a 75%+ decline.",
                 Inches(6.7), Inches(5.75), Inches(5.9), Inches(0.85),
                 font_size=10.5, color=COLOR_LGREY)

    add_footer(slide)
    add_slide_number(slide, 4)
    return slide


# ─────────────────────────────────────────────
# SLIDE 5: IMPORT DRIVERS
# ─────────────────────────────────────────────
def build_slide_05(prs):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, BG_DARK)
    add_rect(slide, Inches(0), Inches(0), Inches(0.08), Inches(7.5), ACCENT_IMPORT)

    add_action_headline(slide,
        "Urea Dependency: Zambia Pays More Each Year to Grow the Same Crops")

    add_text_box(slide, "TOP IMPORTED PRODUCTS (Relative Volume)",
                 Inches(0.3), Inches(1.42), Inches(6.5), Inches(0.28),
                 font_size=8.5, bold=True, color=COLOR_MGREY)

    # Horizontal bar chart for imports
    import_products = [
        ("Urea (Fertilizer)",      100, ACCENT_IMPORT),
        ("Mineral Products",        62, RGBColor(0xf0, 0x7a, 0x30)),
        ("Frozen Goods",            41, RGBColor(0xd4, 0x6a, 0x25)),
        ("Mineral Products-2",      36, RGBColor(0xb8, 0x58, 0x18)),
        ("Ammonia",                 28, RGBColor(0xa0, 0x48, 0x10)),
        ("Other Foods",             22, RGBColor(0x8a, 0x3a, 0x08)),
        ("Crude Petroleum",         19, ACCENT_IMPORT),
        ("Crude Soybean",           14, RGBColor(0xf0, 0x7a, 0x30)),
        ("Palm Oil",                10, RGBColor(0xd4, 0x6a, 0x25)),
        ("Ammonia-2",                8, RGBColor(0xb8, 0x58, 0x18)),
    ]

    bar_left  = Inches(0.3)
    bar_start = Inches(2.2)
    bar_top   = Inches(1.78)
    row_h     = Inches(0.38)
    max_bar_w = Inches(5.5)

    for idx, (name, val, color) in enumerate(import_products):
        by = bar_top + idx * row_h
        bw = max_bar_w * (val / 100)
        add_rect(slide, bar_start, by + Inches(0.04), bw, row_h - Inches(0.09), color)
        add_text_box(slide, name,
                     bar_left, by + Inches(0.07),
                     Inches(1.85), row_h - Inches(0.08),
                     font_size=8.5, color=COLOR_LGREY, align=PP_ALIGN.RIGHT)
        add_text_box(slide, f"{val}",
                     bar_start + bw + Inches(0.07), by + Inches(0.05),
                     Inches(0.5), row_h - Inches(0.08),
                     font_size=8, color=COLOR_MGREY)

    # Urea label
    add_text_box(slide, "#1",
                 bar_start + Inches(0.05), bar_top + Inches(0.08),
                 Inches(0.35), Inches(0.24),
                 font_size=9, bold=True, color=COLOR_WHITE)

    # Right panel: Vicious cycle
    add_rect(slide, Inches(7.1), Inches(1.42), Inches(5.9), Inches(5.65), BG_PANEL)
    add_text_box(slide, "THE FERTILIZER TRAP",
                 Inches(7.3), Inches(1.55), Inches(5.5), Inches(0.3),
                 font_size=10, bold=True, color=ACCENT_IMPORT)

    cycle_steps = [
        (Inches(7.3),  Inches(2.0),  "ZMW falls vs USD", ACCENT_IMPORT),
        (Inches(9.15), Inches(2.65), "Urea costs MORE in ZMW", COLOR_DEFICIT),
        (Inches(7.3),  Inches(3.3),  "Farmers cut usage", COLOR_GOLD),
        (Inches(9.15), Inches(3.95), "Lower crop yields", ACCENT_IMPORT),
        (Inches(7.3),  Inches(4.6),  "Less export revenue", COLOR_DEFICIT),
        (Inches(9.15), Inches(5.25), "More ZMW pressure", COLOR_DEFICIT),
    ]
    arrows = ["↓", "↓", "↓", "↓", "↓"]

    for i, (x, y, text, color) in enumerate(cycle_steps):
        add_rect(slide, x, y, Inches(1.7), Inches(0.5), color)
        add_text_box(slide, text,
                     x + Inches(0.05), y + Inches(0.07),
                     Inches(1.6), Inches(0.38),
                     font_size=8.5, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
        if i < len(cycle_steps) - 1:
            ax = x + Inches(0.85) - Inches(0.08)
            ay = y + Inches(0.5)
            add_text_box(slide, "↓",
                         ax, ay,
                         Inches(0.18), Inches(0.32),
                         font_size=12, bold=True, color=COLOR_MGREY, align=PP_ALIGN.CENTER)

    # Cycle back arrow annotation
    add_text_box(slide, "⟳  Self-reinforcing loop",
                 Inches(7.3), Inches(5.85), Inches(5.5), Inches(0.3),
                 font_size=10, bold=True, color=COLOR_DEFICIT)

    # Key insight
    add_text_box(slide,
                 "\"The import basket is dominated by inputs\nZambia cannot produce domestically.\"",
                 Inches(0.3), Inches(5.82), Inches(6.5), Inches(0.75),
                 font_size=12, italic=True, color=COLOR_GOLD)

    add_footer(slide)
    add_slide_number(slide, 5)
    return slide


# ─────────────────────────────────────────────
# SLIDE 6: EXPORT DRIVERS
# ─────────────────────────────────────────────
def build_slide_06(prs):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, BG_DARK)
    add_rect(slide, Inches(0), Inches(0), Inches(0.08), Inches(7.5), ACCENT_EXPORT)

    add_action_headline(slide,
        "Raw Commodity Exports: High Volume, Low Value-Add, Capped Upside")

    add_text_box(slide, "TOP EXPORTED PRODUCTS (Relative Volume)",
                 Inches(0.3), Inches(1.42), Inches(6.5), Inches(0.28),
                 font_size=8.5, bold=True, color=COLOR_MGREY)

    export_products = [
        ("Tobacco",       100, ACCENT_EXPORT),
        ("Other Raw",      58, RGBColor(0x2a, 0x85, 0xb8)),
        ("Oil Cake",       42, RGBColor(0x1d, 0x6f, 0xa4)),
        ("Maize",          38, RGBColor(0x16, 0x5e, 0x90)),
        ("Cotton",         30, RGBColor(0x10, 0x4e, 0x7c)),
        ("Tobacco-2",      22, ACCENT_EXPORT),
        ("Sweet B.",       15, RGBColor(0x2a, 0x85, 0xb8)),
        ("Maize Seeds",    11, RGBColor(0x1d, 0x6f, 0xa4)),
        ("Raw Ca.",         8, RGBColor(0x16, 0x5e, 0x90)),
        ("Cane Or.",        5, RGBColor(0x10, 0x4e, 0x7c)),
    ]

    bar_left  = Inches(0.3)
    bar_start = Inches(2.2)
    bar_top   = Inches(1.78)
    row_h     = Inches(0.38)
    max_bar_w = Inches(5.5)

    for idx, (name, val, color) in enumerate(export_products):
        by = bar_top + idx * row_h
        bw = max_bar_w * (val / 100)
        add_rect(slide, bar_start, by + Inches(0.04), bw, row_h - Inches(0.09), color)
        add_text_box(slide, name,
                     bar_left, by + Inches(0.07),
                     Inches(1.85), row_h - Inches(0.08),
                     font_size=8.5, color=COLOR_LGREY, align=PP_ALIGN.RIGHT)
        add_text_box(slide, f"{val}",
                     bar_start + bw + Inches(0.07), by + Inches(0.05),
                     Inches(0.5), row_h - Inches(0.08),
                     font_size=8, color=COLOR_MGREY)

    # Right panel
    add_rect(slide, Inches(7.1), Inches(1.42), Inches(5.9), Inches(5.65), BG_PANEL)
    add_text_box(slide, "WHY EXPORTS HAVE A CEILING",
                 Inches(7.3), Inches(1.55), Inches(5.5), Inches(0.3),
                 font_size=10, bold=True, color=ACCENT_EXPORT)

    right_points = [
        (ACCENT_EXPORT,  "PRICE TAKER",
         "Tobacco, maize, cotton are globally\npriced commodities. Zambia has no\nleverage on the USD price."),
        (COLOR_GOLD,     "2021 PEAK → 2022 DIP",
         "Exports hit K14.2bn (2021) but\nFELL to K12.9bn (2022) — first sign\nthat currency tailwind can reverse."),
        (COLOR_DEFICIT,  "STRUCTURAL CAP",
         "Without processing raw commodities\nlocally, export revenue upside is\nlimited regardless of volume."),
    ]

    for idx, (color, title, body) in enumerate(right_points):
        by = Inches(2.0) + idx * Inches(1.72)
        add_rect(slide, Inches(7.3), by, Inches(0.06), Inches(1.38), color)
        add_text_box(slide, title,
                     Inches(7.45), by + Inches(0.03),
                     Inches(5.3), Inches(0.3),
                     font_size=9.5, bold=True, color=color)
        add_text_box(slide, body,
                     Inches(7.45), by + Inches(0.36),
                     Inches(5.3), Inches(1.0),
                     font_size=9.5, color=COLOR_LGREY)

    # Quote
    add_text_box(slide,
                 "\"When commodity prices fall, export revenue falls —\nno matter what the Kwacha does.\"",
                 Inches(0.3), Inches(5.82), Inches(6.5), Inches(0.75),
                 font_size=12, italic=True, color=COLOR_GOLD)

    # Export trend callout
    add_text_box(slide, "2021: K14.17bn  →  2022: K12.95bn  (-K1.22bn)",
                 Inches(7.3), Inches(6.92), Inches(5.5), Inches(0.28),
                 font_size=9, bold=True, color=COLOR_DEFICIT)

    add_footer(slide)
    add_slide_number(slide, 6)
    return slide


# ─────────────────────────────────────────────
# SLIDE 7: THE COMPOUNDING RISK
# ─────────────────────────────────────────────
def build_slide_07(prs):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, BG_DARK)
    add_rect(slide, Inches(0), Inches(0), Inches(0.08), Inches(7.5), COLOR_DEFICIT)

    add_action_headline(slide,
        "The Deficit Loop: How Trade Imbalance Deepens Currency Weakness")

    # Cycle diagram — horizontal flow
    add_text_box(slide, "THE SELF-REINFORCING CYCLE",
                 Inches(0.3), Inches(1.45), Inches(12.5), Inches(0.3),
                 font_size=8.5, bold=True, color=COLOR_MGREY, align=PP_ALIGN.CENTER)

    cycle_nodes = [
        (Inches(0.5),  Inches(2.2), "Trade\nDeficit\nGrows",    COLOR_DEFICIT),
        (Inches(3.0),  Inches(2.2), "Less USD\nInflow to\nZambia",  ACCENT_IMPORT),
        (Inches(5.5),  Inches(2.2), "ZMW\nDepreciates\nFurther",    COLOR_GOLD),
        (Inches(8.0),  Inches(2.2), "Imports Cost\nMore in\nZMW",   ACCENT_IMPORT),
        (Inches(10.5), Inches(2.2), "Deficit\nWidens\nAgain",       COLOR_DEFICIT),
    ]

    node_w = Inches(2.2)
    node_h = Inches(1.35)

    for i, (x, y, text, color) in enumerate(cycle_nodes):
        add_rect(slide, x, y, node_w, node_h, color)
        add_text_box(slide, text,
                     x + Inches(0.05), y + Inches(0.15),
                     node_w - Inches(0.1), node_h - Inches(0.2),
                     font_size=13, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
        if i < len(cycle_nodes) - 1:
            add_text_box(slide, "→",
                         x + node_w + Inches(0.05), y + Inches(0.45),
                         Inches(0.3), Inches(0.5),
                         font_size=22, bold=True, color=COLOR_MGREY, align=PP_ALIGN.CENTER)

    # Loop back arrow label
    add_text_box(slide,
                 "⟲  This loop repeated every year from 2015 to 2022",
                 Inches(0.5), Inches(3.72), Inches(12.2), Inches(0.32),
                 font_size=11, italic=True, color=COLOR_DEFICIT, align=PP_ALIGN.CENTER)

    # Evidence boxes
    add_rect(slide, Inches(0.3), Inches(4.15), Inches(12.5), Inches(0.03), BG_PANEL)

    evidence = [
        ("2013", "SURPLUS\n+K1.2bn", COLOR_GREEN,
         "Exports > Imports.\nSurplus provides USD buffer.\nZMW relatively stable."),
        ("2017", "DEFICIT\n-K4.1bn", ACCENT_IMPORT,
         "First sustained deficit year.\nUrea imports surge.\nZMW starts weakening."),
        ("2020", "DEFICIT\n-K5.8bn", COLOR_DEFICIT,
         "COVID shock. ZMW hits\n~21/USD. Import costs\nexplode in ZMW terms."),
        ("2022", "DEFICIT\n-K9.3bn", COLOR_DEFICIT,
         "Exports FELL despite\nrising ZMW weakness.\nLoop is accelerating."),
    ]

    for idx, (year, status, color, desc) in enumerate(evidence):
        ex = Inches(0.3) + idx * Inches(3.15)
        ey = Inches(4.3)
        add_rect(slide, ex, ey, Inches(2.95), Inches(2.78), BG_PANEL)
        add_rect(slide, ex, ey, Inches(2.95), Inches(0.28), color)
        add_text_box(slide, year,
                     ex + Inches(0.06), ey + Inches(0.03),
                     Inches(1.0), Inches(0.24),
                     font_size=9, bold=True, color=COLOR_WHITE)
        add_text_box(slide, status,
                     ex + Inches(0.06), ey + Inches(0.35),
                     Inches(2.8), Inches(0.65),
                     font_size=18, bold=True, color=color)
        add_text_box(slide, desc,
                     ex + Inches(0.06), ey + Inches(1.1),
                     Inches(2.8), Inches(1.5),
                     font_size=9.5, color=COLOR_LGREY)

    add_text_box(slide,
                 "\"Without intervention, this is a self-reinforcing cycle.\"",
                 Inches(0.3), Inches(7.12), Inches(12.5), Inches(0.28),
                 font_size=10, italic=True, bold=True, color=COLOR_GOLD,
                 align=PP_ALIGN.CENTER)

    add_footer(slide)
    add_slide_number(slide, 7)
    return slide


# ─────────────────────────────────────────────
# SLIDE 8: CALL TO ACTION
# ─────────────────────────────────────────────
def build_slide_08(prs):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, BG_DARK)
    add_rect(slide, Inches(0), Inches(0), Inches(0.08), Inches(7.5), COLOR_GREEN)

    add_action_headline(slide, "Three Priorities to Break the Deficit Loop")

    # Big idea restate
    add_rect(slide, Inches(0.3), Inches(1.42), Inches(12.7), Inches(0.62), BG_PANEL)
    add_text_box(slide,
                 "The K33bn deficit is a structural warning. Currency management alone cannot fix what requires industrial diversification.",
                 Inches(0.45), Inches(1.48), Inches(12.2), Inches(0.52),
                 font_size=12.5, italic=True, bold=True, color=COLOR_GOLD,
                 align=PP_ALIGN.CENTER)

    # Three priority cards
    priorities = [
        (
            "01",
            "Reduce Urea\nDependency",
            ACCENT_IMPORT,
            [
                "Invest in domestic fertilizer production",
                "Explore regional procurement in ZMW or barter",
                "Subsidize alternative/organic inputs",
                "Reduce exposure to USD-priced inputs",
            ]
        ),
        (
            "02",
            "Value-Add\nExports",
            ACCENT_EXPORT,
            [
                "Process tobacco & cotton locally before export",
                "Mill maize into flour for regional markets",
                "Capture more USD per unit exported",
                "Move from price-taker to price-setter",
            ]
        ),
        (
            "03",
            "Import\nSubstitution",
            COLOR_GREEN,
            [
                "Develop local capacity for top import categories",
                "Reduce USD outflow on petroleum & minerals",
                "Industrial policy targeting import-heavy sectors",
                "Regional trade agreements in local currencies",
            ]
        ),
    ]

    card_w = Inches(4.1)
    card_h = Inches(4.38)
    card_top = Inches(2.18)

    for idx, (num, title, color, points) in enumerate(priorities):
        cx = Inches(0.3) + idx * (card_w + Inches(0.15))
        add_rect(slide, cx, card_top, card_w, card_h, BG_PANEL)
        add_rect(slide, cx, card_top, card_w, Inches(0.1), color)

        # Number badge
        add_rect(slide, cx + Inches(0.15), card_top + Inches(0.18),
                 Inches(0.45), Inches(0.45), color)
        add_text_box(slide, num,
                     cx + Inches(0.15), card_top + Inches(0.18),
                     Inches(0.45), Inches(0.45),
                     font_size=12, bold=True, color=COLOR_WHITE,
                     align=PP_ALIGN.CENTER)

        add_text_box(slide, title,
                     cx + Inches(0.72), card_top + Inches(0.2),
                     Inches(3.2), Inches(0.72),
                     font_size=16, bold=True, color=color)

        add_rect(slide, cx + Inches(0.15), card_top + Inches(0.95),
                 card_w - Inches(0.3), Inches(0.02), color)

        for pi, point in enumerate(points):
            py = card_top + Inches(1.08) + pi * Inches(0.72)
            add_rect(slide, cx + Inches(0.18), py + Inches(0.12),
                     Inches(0.12), Inches(0.12), color)
            add_text_box(slide, point,
                         cx + Inches(0.38), py + Inches(0.02),
                         card_w - Inches(0.55), Inches(0.65),
                         font_size=9.5, color=COLOR_LGREY)

    # Bottom bar
    add_rect(slide, Inches(0), Inches(6.78), Inches(13.33), Inches(0.72), BG_PANEL)
    add_text_box(slide,
                 "Data: Zambia Agriculture Import-Export Dashboard 2013–2022  |  FSIO Cohort Analysis  |  May 2026",
                 Inches(0.3), Inches(6.87), Inches(12.7), Inches(0.3),
                 font_size=8.5, color=COLOR_MGREY, align=PP_ALIGN.CENTER)

    add_slide_number(slide, 8)
    return slide


# ─────────────────────────────────────────────
# MAIN: BUILD PRESENTATION
# ─────────────────────────────────────────────
def main():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    print("Building slide 1: Title Slide...")
    build_slide_01(prs)

    print("Building slide 2: The Situation...")
    build_slide_02(prs)

    print("Building slide 3: Diverging Trends...")
    build_slide_03(prs)

    print("Building slide 4: Currency Twist...")
    build_slide_04(prs)

    print("Building slide 5: Import Drivers...")
    build_slide_05(prs)

    print("Building slide 6: Export Drivers...")
    build_slide_06(prs)

    print("Building slide 7: Compounding Risk...")
    build_slide_07(prs)

    print("Building slide 8: Call to Action...")
    build_slide_08(prs)

    output_path = "zambia_trade_presentation.pptx"
    prs.save(output_path)
    print(f"\nPresentation saved: {output_path}")
    print(f"Total slides: {len(prs.slides)}")
    print("Done! Open the .pptx in PowerPoint or LibreOffice Impress.")


if __name__ == "__main__":
    main()
